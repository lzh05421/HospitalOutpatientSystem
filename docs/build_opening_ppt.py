from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_PPT = Path("E:/15751/Desktop/医院门诊挂号与药品管理系统-开题答辩PPT.pptx")
OUT_PPT_COPY = ROOT / "医院门诊挂号与药品管理系统-开题答辩PPT.pptx"
STRUCTURE_IMG = ROOT / "系统功能结构图-重绘版.png"
NOTES = ROOT / "医院门诊挂号与药品管理系统-开题答辩讲稿.md"

TITLE = "基于 Qt 和 MySQL 的医院门诊挂号与药品管理系统"
STUDENT = "学生：刘子航    指导教师：李长毅"
MAJOR = "软件学院 · 计算机科学与技术"

SW, SH = 33.87, 19.05


def rgb(value):
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


NAVY = rgb("17324D")
INK = rgb("1F2937")
MUTED = rgb("667085")
BG = rgb("F4F7FA")
WHITE = rgb("FFFFFF")
ICE = rgb("EAF3F7")
TEAL = rgb("168A8F")
GREEN = rgb("3D8B63")
ORANGE = rgb("D57A35")
RED = rgb("B84A4A")
LINE = rgb("B7C5D5")


def set_text_frame(tf, margin=0.12):
    tf.margin_left = Cm(margin)
    tf.margin_right = Cm(margin)
    tf.margin_top = Cm(margin)
    tf.margin_bottom = Cm(margin)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tf


def text_box(slide, text, x, y, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, font="微软雅黑"):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = set_text_frame(shape.text_frame, 0.04)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return shape


def add_bg(slide, dark=False):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY if dark else BG
    if not dark:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(SW), Cm(0.22))
        band.fill.solid()
        band.fill.fore_color.rgb = TEAL
        band.line.fill.background()


def add_title(slide, title, subtitle=None):
    text_box(slide, title, 1.0, 0.65, 25.8, 0.95, 24, NAVY, True)
    if subtitle:
        text_box(slide, subtitle, 1.05, 1.6, 27.5, 0.55, 9.8, MUTED)
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(29.4), Cm(0.75), Cm(2.35), Cm(0.65))
    tag.fill.solid()
    tag.fill.fore_color.rgb = ICE
    tag.line.color.rgb = rgb("D6E4EC")
    tf = set_text_frame(tag.text_frame, 0.03)
    tf.text = "开题答辩"
    p = tf.paragraphs[0]
    p.font.name = "微软雅黑"
    p.font.size = Pt(9)
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER


def add_footer(slide, idx):
    text_box(slide, f"{idx:02d}", 31.5, 18.05, 1.2, 0.45, 8.5, MUTED, False, PP_ALIGN.RIGHT, "Arial")


def add_card(slide, x, y, w, h, title, body=None, color=TEAL, fill=WHITE, number=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = rgb("D8E3EC")
    shape.line.width = Pt(0.8)
    if number is not None:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x + 0.35), Cm(y + 0.35), Cm(0.88), Cm(0.88))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = set_text_frame(circle.text_frame, 0.0)
        tf.text = str(number)
        p = tf.paragraphs[0]
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        title_x, title_w = x + 1.45, w - 1.75
    else:
        title_x, title_w = x + 0.38, w - 0.76
    text_box(slide, title, title_x, y + 0.28, title_w, 0.55, 12.8, color, True)
    if body:
        box = slide.shapes.add_textbox(Cm(x + 0.38), Cm(y + 1.0), Cm(w - 0.76), Cm(h - 1.2))
        tf = set_text_frame(box.text_frame, 0.03)
        tf.clear()
        for i, item in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.name = "微软雅黑"
            p.font.size = Pt(9.5)
            p.font.color.rgb = rgb("344054")
            p.space_after = Pt(4)
    return shape


def add_stat(slide, x, y, number, label, color):
    text_box(slide, number, x, y, 4.4, 1.1, 30, color, True, PP_ALIGN.CENTER, "Arial")
    text_box(slide, label, x, y + 1.05, 4.4, 0.58, 10, MUTED, False, PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=LINE):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.4)
    return conn


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(13.2), Cm(SH))
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()
    text_box(slide, "基于 Qt 和 MySQL 的\n医院门诊挂号与药品管理系统", 1.15, 2.85, 11.0, 2.9, 24, WHITE, True)
    text_box(slide, "开题答辩 | Qt/C++ · MySQL · C/S 架构", 1.2, 6.05, 10.0, 0.75, 12, rgb("C7DCE9"))
    text_box(slide, MAJOR, 1.2, 13.8, 10.0, 0.55, 11, rgb("D9E8EF"))
    text_box(slide, STUDENT, 1.2, 14.65, 10.8, 0.55, 11, rgb("D9E8EF"))
    for i, (name, x, y, c) in enumerate([
        ("患者预约", 16.0, 3.2, TEAL),
        ("挂号候诊", 22.0, 3.2, GREEN),
        ("医生接诊", 25.5, 7.3, ORANGE),
        ("处方药房", 21.0, 11.4, TEAL),
        ("收费统计", 15.5, 11.4, GREEN),
    ]):
        add_card(slide, x, y, 4.2, 1.45, name, None, c, rgb("F8FBFD"), i + 1)
    add_arrow(slide, 20.2, 3.95, 22.0, 3.95)
    add_arrow(slide, 24.1, 4.65, 25.5, 7.3)
    add_arrow(slide, 25.5, 8.7, 23.2, 11.4)
    add_arrow(slide, 21.0, 12.15, 19.7, 12.15)
    add_footer(slide, 1)


def slide_background(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "课题背景：门诊业务需要数字化闭环", "传统窗口型处理难以支撑多角色协同、实时状态和数据追溯")
    add_stat(slide, 1.0, 3.0, "3", "高频场景", TEAL)
    add_stat(slide, 6.0, 3.0, "3", "关键割裂", ORANGE)
    add_stat(slide, 11.0, 3.0, "1", "完整闭环", GREEN)
    add_card(slide, 1.0, 5.2, 9.8, 4.15, "现实问题", ["患者信息重复录入，窗口处理效率低", "排班、挂号、候诊和收费状态变化快", "药品库存、处方和账单数据容易割裂"], TEAL)
    add_card(slide, 12.0, 5.2, 9.8, 4.15, "建设方向", ["围绕预约诊疗和智慧医院建设要求", "把患者服务、医生诊疗、药房收费一体化", "通过统计和日志提高管理可视化"], GREEN)
    add_card(slide, 23.0, 5.2, 9.0, 4.15, "课题价值", ["覆盖 Qt 客户端、TCP 服务端和 MySQL", "训练需求分析、数据库设计和系统实现能力", "贴近真实医院门诊基础业务形态"], ORANGE)
    add_footer(slide, 2)


def slide_goals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "系统目标与使用对象", "目标不是单一挂号页面，而是门诊挂号、诊疗、药品、收费的基础业务链")
    add_card(slide, 1.0, 2.65, 8.8, 6.5, "建设目标", ["患者可查询排班并预约挂号", "医院人员按角色进入不同业务页面", "排班号源、挂号候诊、处方库存、收费统计形成联动", "服务端统一维护权限、业务规则和数据访问"], TEAL)
    roles = [
        ("患者", "预约挂号", TEAL),
        ("挂号员", "建档/挂号/叫号", GREEN),
        ("医生", "接诊/病历/检查/处方", ORANGE),
        ("药房人员", "审核/发药/库存", TEAL),
        ("收费员", "账单/支付/统计", GREEN),
        ("管理员", "基础数据/日志", NAVY),
    ]
    for i, (role, task, color) in enumerate(roles):
        x = 11.0 + (i % 3) * 7.1
        y = 2.65 + (i // 3) * 3.35
        add_card(slide, x, y, 6.2, 2.45, role, [task], color, WHITE, i + 1)
    add_footer(slide, 3)


def slide_scope(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "功能边界：医院端页面 + 患者预约入口", "源码中已形成菜单权限、页面模块和服务模块的对应关系")
    groups = [
        ("基础与权限", ["院长驾驶舱", "患者管理", "科室管理", "医生管理", "操作日志"], TEAL),
        ("门诊诊疗", ["挂号管理", "候诊队列", "患者病历档案", "医生排班", "医生接诊", "检查检验"], GREEN),
        ("药品收费", ["处方管理", "药品库存", "收费结算", "费用统计"], ORANGE),
    ]
    for i, (title, items, color) in enumerate(groups):
        x = 1.0 + i * 10.8
        add_card(slide, x, 2.55, 9.5, 7.35, title, items, color)
    add_card(slide, 4.6, 11.0, 24.0, 2.7, "权限控制思路", ["管理员拥有全部权限；科主任偏向诊疗与统计；挂号员负责患者、挂号和排班；医生负责接诊、检查和处方；药房负责库存和发药；收费员负责账单和收入统计。"], NAVY, rgb("EEF5F7"))
    add_footer(slide, 4)


def slide_structure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "系统功能结构图", "按照当前项目页面、服务模块和数据库业务链重新整理")
    slide.shapes.add_picture(str(STRUCTURE_IMG), Cm(1.05), Cm(2.45), width=Cm(31.0))
    add_footer(slide, 5)


def slide_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "门诊业务主流程", "从患者预约到收费统计，形成可追溯的闭环链路")
    steps = [
        ("患者建档", "基础信息"),
        ("排班号源", "医生/日期/余号"),
        ("预约挂号", "生成挂号单"),
        ("候诊叫号", "进入医生队列"),
        ("接诊检查", "病历/检查单"),
        ("处方发药", "审核/出库"),
        ("收费结算", "账单/支付"),
        ("统计日志", "收入/追溯"),
    ]
    x0, y0, w, h, gap = 1.0, 4.0, 3.55, 2.3, 0.48
    for i, (name, desc) in enumerate(steps):
        x = x0 + i * (w + gap)
        add_card(slide, x, y0, w, h, name, [desc], [TEAL, GREEN, ORANGE, NAVY][i % 4], WHITE, i + 1)
        if i < len(steps) - 1:
            add_arrow(slide, x + w, y0 + h / 2, x + w + gap - 0.05, y0 + h / 2)
    lanes = [
        ("患者/挂号员", "患者信息、挂号单、预约时段"),
        ("医生/检查人员", "候诊队列、病历记录、检查结果"),
        ("药房/收费员", "处方状态、库存出库、账单支付"),
        ("管理员/科主任", "科室医生、排班规则、统计看板、日志审计"),
    ]
    for i, (role, data) in enumerate(lanes):
        y = 9.0 + i * 1.35
        add_card(slide, 2.0, y, 7.0, 0.9, role, None, TEAL, rgb("EEF7F8"))
        text_box(slide, data, 9.6, y + 0.1, 22.0, 0.6, 11, INK)
    add_footer(slide, 6)


def slide_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "技术架构：Qt/C++ 客户端 + Qt TCP 服务端 + MySQL", "客户端专注交互，服务端统一业务规则，数据库集中保存数据")
    layers = [
        ("客户端表现层", ["Qt Widgets 桌面界面", "患者预约入口", "角色菜单与表格页面"], TEAL),
        ("网络通信层", ["QTcpSocket", "TCP + JSON Lines", "统一请求/响应结构"], GREEN),
        ("服务端业务层", ["HospitalServer", "RequestRouter 鉴权", "业务 Service 分发处理"], ORANGE),
        ("数据持久层", ["MySQL", "QODBC/QMYSQL", "演示模式兜底"], NAVY),
    ]
    x = 1.0
    for i, (title, items, color) in enumerate(layers):
        add_card(slide, x, 3.0, 7.1, 6.2, title, items, color)
        if i < len(layers) - 1:
            add_arrow(slide, x + 7.1, 6.1, x + 8.1, 6.1)
        x += 8.1
    add_card(slide, 3.2, 11.2, 27.3, 2.4, "设计要点", ["客户端不直接操作数据库，服务端负责权限校验和业务一致性；数据库集中保存患者、医生、排班、挂号、病历、处方、药品、账单、支付、统计和日志数据。"], TEAL, rgb("EEF5F7"))
    add_footer(slide, 7)


def slide_database(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "数据库设计", "围绕角色、基础资料、诊疗业务、药品收费四组核心数据建模")
    groups = [
        ("账号权限", ["roles", "users", "operation_logs", "audit_log_details"], TEAL),
        ("基础资料", ["departments", "doctors", "patients", "doctor_schedules"], GREEN),
        ("诊疗业务", ["registrations", "medical_records", "examinations", "prescriptions"], ORANGE),
        ("药品收费", ["drug_categories", "drugs", "prescription_items", "stock_records", "bills", "payments", "fee_statistics_daily"], NAVY),
    ]
    x = 1.0
    for i, (title, tables, color) in enumerate(groups):
        add_card(slide, x, 2.7, 7.4, 6.7, title, tables, color)
        if i < len(groups) - 1:
            add_arrow(slide, x + 7.4, 5.95, x + 8.2, 5.95)
        x += 8.2
    notes = [
        "患者/医生/排班生成挂号记录",
        "挂号关联病历、检查和处方",
        "处方关联药品库存，账单和支付完成收费闭环",
    ]
    for i, note in enumerate(notes):
        add_card(slide, 2.0 + i * 10.6, 11.0, 9.2, 1.55, note, None, [TEAL, GREEN, ORANGE][i], rgb("F8FBFD"), i + 1)
    add_footer(slide, 8)


def slide_implementation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "核心实现方案", "用项目中的真实工程结构支撑答辩")
    cards = [
        ("通信协议", ["QTcpSocket 长连接", "JSON Lines 请求/响应", "module + action 定位接口"], TEAL),
        ("权限路由", ["RequestRouter 统一鉴权", "按角色控制模块动作", "公开接口仅限排班查询和挂号创建"], GREEN),
        ("业务服务", ["Patient/Registration/Schedule", "Consultation/Examination", "Prescription/Inventory/Billing"], ORANGE),
        ("数据访问", ["DatabaseManager 连接 MySQL", "QODBC/QMYSQL", "连接异常提示与演示模式"], NAVY),
    ]
    for i, (title, items, color) in enumerate(cards):
        add_card(slide, 1.0 + i * 8.1, 2.8, 7.1, 5.4, title, items, color, WHITE, i + 1)
    text_box(slide, "ApiClient → Protocol → HospitalServer → RequestRouter → Service → DatabaseManager → MySQL", 2.0, 10.3, 29.6, 1.1, 16, NAVY, True, PP_ALIGN.CENTER, "Consolas")
    add_card(slide, 4.2, 12.2, 25.4, 2.1, "答辩强调点", ["系统不是固定演示数据：患者预约、医生排班、药品入库等核心写入均可落到 MySQL；页面支持查询、分页、刷新和状态维护。"], TEAL, rgb("EEF5F7"))
    add_footer(slide, 9)


def slide_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "工作方案与进度安排", "按软件工程过程推进，每个阶段都有明确产出")
    periods = [
        ("需求与开题", "2025.12", "调研/需求/模块"),
        ("总体设计", "2026.01", "架构/数据库/原型"),
        ("核心开发", "2026.02-03", "挂号/排班/接诊"),
        ("扩展开发", "2026.03-04", "检查/处方/收费"),
        ("测试部署", "2026.04-05", "联调、多环境验证"),
        ("论文答辩", "2026.05-06", "论文修改、PPT、提交"),
    ]
    for i, (name, time, desc) in enumerate(periods):
        x = 1.0 + i * 5.35
        add_card(slide, x, 4.1, 4.65, 4.6, name, [time, desc], [TEAL, GREEN, ORANGE, NAVY, RED, TEAL][i], WHITE, i + 1)
        if i < len(periods) - 1:
            add_arrow(slide, x + 4.65, 6.4, x + 5.25, 6.4)
    add_card(slide, 3.4, 11.0, 27.0, 2.7, "测试重点", ["按患者、挂号员、医生、药房人员、收费员、科主任、管理员等角色分别验证；重点检查号源扣减、候诊状态、处方发药、库存变化、收费统计和操作日志。"], GREEN, rgb("EEF5F7"))
    add_footer(slide, 10)


def slide_outcomes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "预期成果与项目特色", "答辩时突出完整业务链和工程化实现")
    add_card(slide, 1.0, 2.65, 9.6, 6.6, "预期成果", ["可运行的 C/S 门诊管理系统", "MySQL 建库脚本和演示数据", "跨平台构建与运行脚本", "开题报告、毕业论文、答辩 PPT"], TEAL)
    add_card(slide, 12.0, 2.65, 9.6, 6.6, "项目特色", ["患者入口与医院人员入口分离", "角色权限驱动菜单和接口", "排班号源影响预约挂号", "处方、药品、收费形成联动"], GREEN)
    add_card(slide, 23.0, 2.65, 9.0, 6.6, "可展示亮点", ["候诊队列与叫号状态", "扫码入库与库存预警", "院长驾驶舱与费用统计", "操作日志与审计追踪"], ORANGE)
    add_footer(slide, 11)


def slide_risks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "风险分析与解决措施", "提前准备答辩老师可能关注的问题")
    risks = [
        ("数据库环境差异", ["配置文件管理连接参数", "支持 QODBC/QMYSQL 两种方向", "保留演示模式验证界面和流程"], RED),
        ("业务状态不同步", ["挂号成功扣减号源", "列表支持手动刷新和自动刷新", "核心状态由服务端统一维护"], ORANGE),
        ("模块多导致测试遗漏", ["按角色设计测试用例", "从建档到收费做端到端验证", "操作日志辅助定位问题"], TEAL),
    ]
    for i, (title, items, color) in enumerate(risks):
        add_card(slide, 1.2 + i * 10.9, 3.0, 9.4, 7.4, title, items, color, WHITE, i + 1)
    add_card(slide, 4.0, 12.2, 25.8, 2.2, "答辩准备", ["现场演示优先走主流程：患者预约 → 医院端挂号管理 → 候诊叫号 → 医生接诊 → 处方/药品 → 收费统计。"], NAVY, rgb("EEF5F7"))
    add_footer(slide, 12)


def closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, True)
    text_box(slide, "请各位老师批评指正", 5.0, 6.3, 23.8, 1.35, 34, WHITE, True, PP_ALIGN.CENTER)
    text_box(slide, TITLE, 5.0, 8.2, 23.8, 0.8, 14, rgb("C7DCE9"), False, PP_ALIGN.CENTER)
    text_box(slide, STUDENT, 5.0, 9.45, 23.8, 0.65, 11, rgb("D9E8EF"), False, PP_ALIGN.CENTER)
    add_footer(slide, 13)


def build_ppt():
    if not STRUCTURE_IMG.exists():
        raise FileNotFoundError(f"缺少结构图：{STRUCTURE_IMG}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    slide_background(prs)
    slide_goals(prs)
    slide_scope(prs)
    slide_structure(prs)
    slide_flow(prs)
    slide_arch(prs)
    slide_database(prs)
    slide_implementation(prs)
    slide_plan(prs)
    slide_outcomes(prs)
    slide_risks(prs)
    closing(prs)
    prs.save(str(OUT_PPT))
    shutil.copy2(OUT_PPT, OUT_PPT_COPY)
    return OUT_PPT


def build_notes():
    NOTES.write_text(
        """# 医院门诊挂号与药品管理系统开题答辩讲稿

1. 开场先说明课题：系统基于 Qt/C++、Qt TCP 服务端和 MySQL，面向门诊挂号、诊疗、药品、收费等基础业务。
2. 背景页强调门诊业务高频、状态变化快，传统表格或纸质单据容易造成重复录入、号源不同步和统计困难。
3. 目标页说明使用对象：患者、挂号员、医生、药房人员、收费员、管理员和科主任，不同角色看到不同功能。
4. 功能页用 12 个模块概括：登录权限、驾驶舱、患者病历、科室医生、排班号源、挂号候诊、诊疗检查、处方药房、药品库存、收费结算、统计分析、系统运维。
5. 技术架构页重点讲 C/S：客户端 ApiClient 发送 JSON Lines 请求，服务端 RequestRouter 鉴权并分发到业务 Service，DatabaseManager 访问 MySQL。
6. 数据库页说明主链路：患者、医生、排班生成挂号；挂号关联病历、检查、处方；处方关联药品库存；账单和支付完成收费闭环。
7. 实施计划页按需求、设计、核心开发、扩展开发、测试部署、论文答辩展开。
8. 预期成果页突出可运行系统、数据库脚本、跨平台构建脚本、论文和答辩材料。
9. 风险页准备回答数据库环境、业务状态同步、模块测试覆盖等问题。
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    output = build_ppt()
    build_notes()
    print(output)
    print(OUT_PPT_COPY)
    print(NOTES)

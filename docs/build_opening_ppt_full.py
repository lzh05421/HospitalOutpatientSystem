from pathlib import Path
import importlib.util
import shutil

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches


ROOT = Path(__file__).resolve().parent
BASE_PATH = ROOT / "build_opening_ppt_hq.py"
OUT_PPT = Path("E:/15751/Desktop/医院门诊挂号与药品管理系统-开题答辩PPT-全功能版.pptx")
OUT_PPT_ALT = Path("E:/15751/Desktop/医院门诊挂号与药品管理系统-开题答辩PPT-全功能版-新版.pptx")
OUT_COPY = ROOT / "医院门诊挂号与药品管理系统-开题答辩PPT-全功能版.pptx"
OUT_STRUCTURE = ROOT / "系统功能结构图-全功能版.png"
OUT_NOTES = ROOT / "开题答辩讲稿-全功能版.md"


spec = importlib.util.spec_from_file_location("ppt_hq", BASE_PATH)
base = importlib.util.module_from_spec(spec)
spec.loader.exec_module(base)


FULL_MODULES = [
    ("入口与登录", ["患者预约入口", "医院人员入口", "角色预选", "账号密码登录", "连接状态提示", "会话信息保存"]),
    ("权限与日志", ["角色菜单控制", "接口动作授权", "公开预约接口", "写操作日志", "审计明细记录"]),
    ("患者预约挂号", ["门诊大类选择", "专科诊室选择", "可约号源刷新", "医生日期时段", "症状辅助选科", "就诊人信息", "预约记录显示"]),
    ("院长驾驶舱", ["今日挂号", "候诊人数", "已完成接诊", "今日收入", "收入构成", "库存预警", "医生接诊排行"]),
    ("患者管理", ["患者列表", "患者身份查询", "信息修改", "删除保护", "就诊次数", "复诊提示", "档案完整度"]),
    ("挂号管理", ["挂号记录", "科室医生筛选", "日期状态查询", "挂号状态维护", "退号/取消", "自动刷新", "号源扣减"]),
    ("候诊队列", ["候诊列表", "科室筛选", "医生筛选", "状态筛选", "排队序号", "预计等待", "叫号操作"]),
    ("患者病历档案", ["患者基础信息", "挂号记录", "接诊病历", "检查记录", "处方账单", "档案修改", "历史追溯"]),
    ("科室管理", ["科室编码", "科室名称", "位置维护", "新增科室", "修改科室", "停用科室"]),
    ("医生管理", ["医生新增", "医生修改", "医生停用", "科室联动", "职称专长", "联系电话", "挂号费"]),
    ("医生排班", ["医生列表刷新", "门诊级联选择", "排班保存", "号源调整", "停诊排班", "智能排班7天", "长期规则", "未排班医生", "重新排班"]),
    ("医生接诊", ["候诊查看", "医生叫号", "开始接诊", "主诉诊断医嘱", "辅助诊断提示", "申请检查", "开药明细", "检查后候诊", "保存病历"]),
    ("检查检验", ["检查单列表", "开检查单", "项目说明", "结果录入", "完成状态", "回到候诊"]),
    ("处方管理", ["处方列表", "开立处方", "药品明细", "用法频次天数", "过敏与患者情况", "处方金额", "处方审核", "确认发药", "库存扣减"]),
    ("药品库存", ["药品列表", "条形码", "药品分类", "规格单位", "进价售价", "扫码入库", "新药入库", "库存预警", "出入库记录", "有效期/管控"]),
    ("收费结算", ["账单列表", "费用合计", "确认收费", "支付方式", "支付记录", "退费处理", "状态维护"]),
    ("费用统计", ["日收入统计", "科室统计", "挂号收入", "药品收入", "图表展示", "CSV报表导出"]),
    ("通用列表能力", ["关键字查询", "分组筛选", "医生筛选", "号别筛选", "日期筛选", "分页显示", "手动刷新", "自动刷新", "CSV导出"]),
    ("系统配置部署", ["CMake工程", "Qt Creator运行", "VS Code任务", "Windows脚本", "Linux/麒麟脚本", "MySQL配置", "ODBC/QMYSQL", "演示模式兜底"]),
]


def wrapped(draw, text, font, max_width):
    lines = []
    cur = ""
    for ch in text:
        test = cur + ch
        box = draw.textbbox((0, 0), test, font=font)
        if box[2] - box[0] <= max_width or not cur:
            cur = test
        else:
            lines.append(cur)
            cur = ch
    if cur:
        lines.append(cur)
    return lines


def draw_full_structure(path):
    width, height = 4200, 2400
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)
    title_font = base.pil_font(54, True)
    root_font = base.pil_font(40, True)
    module_font = base.pil_font(30, True)
    item_font = base.pil_font(23)
    small_font = base.pil_font(21)
    caption_font = base.pil_font(32)

    navy = (20, 50, 74)
    teal = (22, 138, 143)
    green = (60, 140, 99)
    amber = (215, 130, 58)
    red = (184, 77, 77)
    border = (48, 66, 84)
    soft = (244, 248, 250)
    colors = [teal, green, amber, navy, red]

    root = (1230, 42, 2970, 120)
    draw.rectangle(root, fill=soft, outline=border, width=3)
    base.draw_center(draw, root, "基于 Qt 和 MySQL 的医院门诊挂号与药品管理系统", title_font, navy)

    subtitle = f"全功能结构图：{len(FULL_MODULES)} 个模块，{sum(len(items) for _, items in FULL_MODULES)} 个功能点"
    base.draw_center(draw, (0, 145, width, 190), subtitle, root_font, teal)

    cols = 4
    card_w, card_h = 930, 335
    gap_x, gap_y = 70, 54
    start_x, start_y = 135, 260

    root_cx = (root[0] + root[2]) // 2
    trunk_y = 220
    draw.line((root_cx, root[3], root_cx, trunk_y), fill=border, width=4)
    draw.line((start_x + card_w // 2, trunk_y, start_x + (cols - 1) * (card_w + gap_x) + card_w // 2, trunk_y), fill=border, width=4)

    for idx, (name, items) in enumerate(FULL_MODULES):
        row, col = divmod(idx, cols)
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        color = colors[idx % len(colors)]
        cx = x + card_w // 2
        if row == 0:
            draw.line((cx, trunk_y, cx, y), fill=border, width=3)
        draw.rounded_rectangle((x, y, x + card_w, y + card_h), radius=18, fill=(250, 253, 255), outline=border, width=3)
        draw.rounded_rectangle((x, y, x + card_w, y + 62), radius=18, fill=(226, 241, 240), outline=(226, 241, 240), width=0)
        base.draw_center(draw, (x + 20, y + 8, x + card_w - 20, y + 58), name, module_font, color)

        col_w = (card_w - 86) // 2
        left_x = x + 36
        right_x = x + 36 + col_w + 34
        line_h = 34 if len(items) <= 8 else 29
        font = item_font if len(items) <= 8 else small_font
        for i, item in enumerate(items):
            list_col = 0 if i < (len(items) + 1) // 2 else 1
            row_i = i if list_col == 0 else i - (len(items) + 1) // 2
            tx = left_x if list_col == 0 else right_x
            ty = y + 88 + row_i * line_h
            draw.ellipse((tx, ty + 8, tx + 10, ty + 18), fill=color)
            draw.text((tx + 20, ty), item, font=font, fill=(31, 43, 60))

    base.draw_center(draw, (0, 2290, width, 2350), "图1 系统功能结构图（全功能版）", caption_font, (20, 32, 46))
    img.save(path)


def scope_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    base.bg(s)
    base.title(s, "功能范围：按项目真实页面和服务重新梳理", "不再固定 12 个模块，也不强行每个模块 3 个子功能")
    groups = [
        ("入口与基础", ["入口与登录", "权限与日志", "院长驾驶舱", "通用列表能力"], base.TEAL),
        ("门诊主流程", ["患者预约挂号", "患者管理", "挂号管理", "候诊队列", "患者病历档案"], base.GREEN),
        ("诊疗与药房", ["科室管理", "医生管理", "医生排班", "医生接诊", "检查检验", "处方管理", "药品库存"], base.AMBER),
        ("收费与运维", ["收费结算", "费用统计", "系统配置部署"], base.NAVY),
    ]
    for i, (h, items, c) in enumerate(groups):
        x = 1.0 + (i % 2) * 16.0
        y = 2.35 + (i // 2) * 5.15
        base.card(s, x, y, 14.2, 4.65, h, items, c, base.WHITE, body_size=8.9)
    base.card(s, 4.2, 13.35, 25.4, 1.45, "全功能版口径", [f"结构图包含 {len(FULL_MODULES)} 个模块、{sum(len(items) for _, items in FULL_MODULES)} 个功能点，来源于 MainWindow 菜单、页面按钮、服务端 action 和数据库表。"], base.TEAL, base.rgb("EEF7F8"), body_size=8.8)
    base.footer(s, 4)


def structure_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    base.bg(s)
    base.title(s, "系统功能结构图（全功能版）", "按当前项目所有主要页面、服务动作和通用能力展开")
    s.shapes.add_picture(str(OUT_STRUCTURE), Cm(0.75), Cm(2.1), width=Cm(32.3))
    base.footer(s, 5)


def build_ppt():
    draw_full_structure(OUT_STRUCTURE)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    base.cover(prs)
    base.background_slide(prs)
    base.goal_slide(prs)
    scope_slide(prs)
    structure_slide(prs)
    base.flow_slide(prs)
    base.architecture_slide(prs)
    base.database_slide(prs)
    base.implementation_slide(prs)
    base.demo_slide(prs)
    base.plan_slide(prs)
    base.outcomes_slide(prs)
    base.risk_slide(prs)
    base.close_slide(prs)
    prs.save(str(OUT_COPY))
    try:
        shutil.copy2(OUT_COPY, OUT_PPT)
    except PermissionError:
        shutil.copy2(OUT_COPY, OUT_PPT_ALT)


def build_notes():
    OUT_NOTES.write_text(
        f"""# 开题答辩讲稿提纲（全功能版）

1. 说明这版结构图按当前项目源码重新梳理，包含 {len(FULL_MODULES)} 个模块、{sum(len(items) for _, items in FULL_MODULES)} 个功能点。
2. 入口与登录：患者预约入口、医院人员入口、角色预选、账号密码登录、会话保存。
3. 门诊主流程：患者预约挂号、挂号管理、候诊队列、医生接诊、检查检验、处方管理、收费结算。
4. 基础数据：患者管理、患者病历档案、科室管理、医生管理、医生排班。
5. 药品与收费：药品库存、处方审核发药、账单收费、退费、费用统计、CSV 报表。
6. 管理支撑：院长驾驶舱、权限控制、操作日志、审计明细、通用查询筛选、分页、自动刷新、导出。
7. 运维支撑：CMake 工程、Qt Creator/VS Code、Windows/Linux/银河麒麟脚本、MySQL 配置、演示模式兜底。
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    build_ppt()
    build_notes()
    print(OUT_PPT)
    print(OUT_COPY)
    print(OUT_STRUCTURE)
    print(OUT_NOTES)

from pathlib import Path
import shutil

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_PPT = Path("E:/15751/Desktop/医院门诊挂号与药品管理系统-开题答辩PPT-高质量版.pptx")
OUT_COPY = ROOT / "医院门诊挂号与药品管理系统-开题答辩PPT-高质量版.pptx"
OUT_STRUCTURE = ROOT / "系统功能结构图-树状高质量版.png"
OUT_NOTES = ROOT / "开题答辩讲稿-高质量版.md"

TITLE = "基于 Qt 和 MySQL 的医院门诊挂号与药品管理系统"
STUDENT = "学生：刘子航    指导教师：李长毅"
MAJOR = "软件学院 · 计算机科学与技术"

SW, SH = 33.87, 19.05


def rgb(value):
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


NAVY = rgb("14324A")
TEAL = rgb("168A8F")
MINT = rgb("DCEFEB")
GREEN = rgb("3C8C63")
AMBER = rgb("D7823A")
RED = rgb("B84D4D")
INK = rgb("1F2937")
MUTED = rgb("667085")
BG = rgb("F4F8FA")
WHITE = rgb("FFFFFF")
LINE = rgb("B7C8D6")


def pil_font(size, bold=False):
    candidates = [
        Path(r"C:\Windows\Fonts\simhei.ttf") if bold else Path(r"C:\Windows\Fonts\simsun.ttc"),
        Path(r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
    ]
    for item in candidates:
        if item.exists():
            return ImageFont.truetype(str(item), size)
    return ImageFont.load_default()


def draw_center(draw, box, text, font_obj, fill=(20, 32, 46)):
    x1, y1, x2, y2 = box
    bbox = draw.multiline_textbbox((0, 0), text, font=font_obj, spacing=4)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.multiline_text(
        (x1 + (x2 - x1 - tw) / 2, y1 + (y2 - y1 - th) / 2 - 2),
        text,
        font=font_obj,
        fill=fill,
        spacing=4,
        align="center",
    )


def draw_vertical(draw, box, text, font_obj, fill=(20, 32, 46)):
    x1, y1, x2, y2 = box
    chars = list(text)
    line_h = 36
    total_h = line_h * len(chars)
    y = y1 + max(8, (y2 - y1 - total_h) / 2)
    for ch in chars:
        bbox = draw.textbbox((0, 0), ch, font=font_obj)
        tw = bbox[2] - bbox[0]
        draw.text((x1 + (x2 - x1 - tw) / 2, y), ch, font=font_obj, fill=fill)
        y += line_h


def draw_tree_structure(path):
    width, height = 3600, 1650
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)
    root_font = pil_font(50, True)
    module_font = pil_font(34, True)
    leaf_font = pil_font(30)
    caption_font = pil_font(34)
    border = (28, 42, 58)
    line = (28, 42, 58)
    root_fill = (242, 247, 250)
    module_fill = (236, 246, 244)
    leaf_fill = (255, 255, 255)

    modules = [
        ("登录权限", ["医院人员登录", "角色权限控制", "操作日志管理"]),
        ("患者管理", ["患者建档", "患者查询", "病历档案"]),
        ("挂号管理", ["科室医生", "日期时段", "号源控制"]),
        ("候诊叫号", ["候诊队列", "科室筛选", "医生叫号"]),
        ("医生管理", ["医生信息", "科室维护", "状态管理"]),
        ("医生排班", ["排班维护", "号源调整", "停诊管理"]),
        ("医生接诊", ["待诊查看", "病历填写", "检查候诊"]),
        ("检查检验", ["开检查单", "结果录入", "状态查询"]),
        ("处方管理", ["开立处方", "处方审核", "确认发药"]),
        ("药品库存", ["药品维护", "扫码入库", "库存预警"]),
        ("收费结算", ["账单查询", "缴费退费", "支付记录"]),
        ("费用统计", ["日收入统计", "科室统计", "图表展示"]),
    ]

    root = (1330, 50, 2270, 142)
    draw.rectangle(root, fill=root_fill, outline=border, width=3)
    draw_center(draw, root, "基于 Qt 和 MySQL 的医院\n挂号与药品管理系统", root_font, (10, 31, 51))

    module_w = 188
    module_h = 62
    gap = 84
    total = len(modules) * module_w + (len(modules) - 1) * gap
    start_x = (width - total) // 2
    trunk_y = 280
    mod_y = 360
    leaf_y1, leaf_y2 = 560, 1340
    leaf_w, leaf_gap = 56, 18

    centers = []
    for i in range(len(modules)):
        x = start_x + i * (module_w + gap)
        centers.append(x + module_w // 2)

    root_cx = (root[0] + root[2]) // 2
    draw.line((root_cx, root[3], root_cx, trunk_y), fill=line, width=3)
    draw.line((centers[0], trunk_y, centers[-1], trunk_y), fill=line, width=3)

    for i, (name, leaves) in enumerate(modules):
        x = start_x + i * (module_w + gap)
        cx = centers[i]
        module_box = (x, mod_y, x + module_w, mod_y + module_h)
        draw.line((cx, trunk_y, cx, mod_y), fill=line, width=3)
        draw.rectangle(module_box, fill=module_fill, outline=border, width=3)
        draw_center(draw, module_box, name, module_font, (14, 83, 82))

        leaf_total = 3 * leaf_w + 2 * leaf_gap
        leaf_start = cx - leaf_total // 2
        leaf_centers = []
        branch_y = 492
        draw.line((cx, mod_y + module_h, cx, branch_y), fill=line, width=3)
        for j, leaf in enumerate(leaves):
            lx = leaf_start + j * (leaf_w + leaf_gap)
            leaf_centers.append(lx + leaf_w // 2)
        draw.line((leaf_centers[0], branch_y, leaf_centers[-1], branch_y), fill=line, width=3)
        for j, leaf in enumerate(leaves):
            lx = leaf_start + j * (leaf_w + leaf_gap)
            lc = lx + leaf_w // 2
            draw.line((lc, branch_y, lc, leaf_y1), fill=line, width=3)
            leaf_box = (lx, leaf_y1, lx + leaf_w, leaf_y2)
            draw.rectangle(leaf_box, fill=leaf_fill, outline=border, width=3)
            draw_vertical(draw, leaf_box, leaf, leaf_font, (13, 31, 47))

    draw_center(draw, (0, 1455, width, 1515), "图1 系统功能结构图", caption_font, (20, 32, 46))
    img.save(path)


def set_tf(tf, margin=0.12):
    tf.margin_left = Cm(margin)
    tf.margin_right = Cm(margin)
    tf.margin_top = Cm(margin)
    tf.margin_bottom = Cm(margin)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tf


def txt(slide, text, x, y, w, h, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT, font="微软雅黑"):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = set_tf(box.text_frame, 0.04)
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def bg(slide, dark=False):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY if dark else BG
    if not dark:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(SW), Cm(0.22))
        band.fill.solid()
        band.fill.fore_color.rgb = TEAL
        band.line.fill.background()


def title(slide, text, subtitle=None):
    txt(slide, text, 1.0, 0.62, 27.0, 0.9, 24, NAVY, True)
    if subtitle:
        txt(slide, subtitle, 1.04, 1.55, 28.2, 0.55, 9.8, MUTED)
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(29.2), Cm(0.7), Cm(2.6), Cm(0.65))
    tag.fill.solid()
    tag.fill.fore_color.rgb = rgb("E7F2F2")
    tag.line.color.rgb = rgb("C9E0E0")
    tf = set_tf(tag.text_frame, 0.02)
    tf.text = "开题答辩"
    p = tf.paragraphs[0]
    p.font.name = "微软雅黑"
    p.font.size = Pt(9)
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER


def footer(slide, n):
    txt(slide, f"{n:02d}", 31.4, 18.05, 1.1, 0.45, 8.5, MUTED, False, PP_ALIGN.RIGHT, "Arial")


def card(slide, x, y, w, h, heading, body=None, color=TEAL, fill=WHITE, num=None, heading_size=12.5, body_size=9.4):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = rgb("D6E2EA")
    s.line.width = Pt(0.85)
    if num is not None:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x + 0.34), Cm(y + 0.34), Cm(0.86), Cm(0.86))
        c.fill.solid()
        c.fill.fore_color.rgb = color
        c.line.fill.background()
        tf = set_tf(c.text_frame, 0.0)
        tf.text = str(num)
        p = tf.paragraphs[0]
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        hx, hw = x + 1.42, w - 1.7
    else:
        hx, hw = x + 0.42, w - 0.84
    txt(slide, heading, hx, y + 0.26, hw, 0.55, heading_size, color, True)
    if body:
        box = slide.shapes.add_textbox(Cm(x + 0.42), Cm(y + 0.96), Cm(w - 0.84), Cm(h - 1.15))
        tf = set_tf(box.text_frame, 0.03)
        tf.clear()
        for i, item in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.name = "微软雅黑"
            p.font.size = Pt(body_size)
            p.font.color.rgb = rgb("344054")
            p.space_after = Pt(3.5)
    return s


def line(slide, x1, y1, x2, y2, color=LINE, width=1.25):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, True)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(13.2), Cm(SH))
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()
    txt(s, "基于 Qt 和 MySQL 的\n医院门诊挂号与药品管理系统", 1.1, 2.85, 11.2, 3.0, 24, WHITE, True)
    txt(s, "开题答辩 | Qt/C++ · MySQL · C/S 架构", 1.15, 6.25, 10.2, 0.7, 12, rgb("C9DDE8"))
    txt(s, MAJOR, 1.15, 13.7, 10.2, 0.55, 11, rgb("DCE9EF"))
    txt(s, STUDENT, 1.15, 14.55, 11.0, 0.55, 11, rgb("DCE9EF"))
    steps = [("患者预约", 16.0, 3.0, TEAL), ("挂号候诊", 22.2, 3.0, GREEN), ("医生接诊", 25.7, 7.2, AMBER), ("处方药房", 21.0, 11.3, TEAL), ("收费统计", 15.8, 11.3, GREEN)]
    for i, (name, x, y, c) in enumerate(steps, 1):
        card(s, x, y, 4.25, 1.4, name, None, c, rgb("F8FBFD"), i)
    line(s, 20.25, 3.7, 22.2, 3.7)
    line(s, 24.35, 4.25, 25.7, 7.2)
    line(s, 25.7, 8.55, 23.15, 11.3)
    line(s, 21.0, 12.0, 20.05, 12.0)
    footer(s, 1)


def background_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "课题背景：门诊业务需要一个数据闭环", "挂号、候诊、诊疗、药品、收费和统计之间需要实时同步")
    for x, number, label, c in [(1.1, "3", "高频环节", TEAL), (6.0, "6", "角色协同", AMBER), (10.9, "1", "业务闭环", GREEN)]:
        txt(s, number, x, 2.6, 3.7, 1.0, 30, c, True, PP_ALIGN.CENTER, "Arial")
        txt(s, label, x, 3.58, 3.7, 0.45, 9.5, MUTED, False, PP_ALIGN.CENTER)
    card(s, 1.0, 5.2, 9.5, 4.0, "现实问题", ["纸质单据或表格容易重复录入", "号源、候诊、收费状态变化快", "药品库存、处方和账单容易割裂"], TEAL)
    card(s, 12.0, 5.2, 9.5, 4.0, "系统价值", ["把患者服务、医生诊疗、药房收费连接起来", "通过统一数据库保证数据一致", "操作日志和统计图表支撑管理追溯"], GREEN)
    card(s, 23.0, 5.2, 8.9, 4.0, "毕设训练点", ["Qt 桌面端", "TCP 服务端", "MySQL 数据库", "跨平台构建和测试"], AMBER)
    footer(s, 2)


def goal_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "系统目标与用户角色", "患者入口和医院人员入口分离，内部人员按角色使用不同业务模块")
    card(s, 1.0, 2.55, 8.9, 6.7, "建设目标", ["患者可查询排班并预约挂号", "医院人员通过账号密码登录", "服务端统一完成权限校验和业务处理", "MySQL 保存患者、医生、挂号、处方、库存、账单等数据"], TEAL)
    roles = [("患者", "预约挂号", TEAL), ("挂号员", "建档/挂号/叫号", GREEN), ("医生", "接诊/检查/处方", AMBER), ("药房人员", "审核/发药/库存", TEAL), ("收费员", "账单/支付/统计", GREEN), ("管理员", "基础数据/日志", NAVY)]
    for i, (r, task, c) in enumerate(roles):
        x = 11.1 + (i % 3) * 7.1
        y = 2.55 + (i // 3) * 3.35
        card(s, x, y, 6.2, 2.45, r, [task], c, WHITE, i + 1)
    footer(s, 3)


def scope_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "功能范围：12 个主模块，36 个子功能", "按你给出的结构图口径，重新统一开题报告和答辩内容")
    groups = [
        ("基础信息与权限", ["登录权限", "患者管理", "医生管理"], TEAL),
        ("挂号与诊疗流程", ["挂号管理", "候诊叫号", "医生排班", "医生接诊", "检查检验"], GREEN),
        ("药品收费与统计", ["处方管理", "药品库存", "收费结算", "费用统计"], AMBER),
    ]
    for i, (h, items, c) in enumerate(groups):
        card(s, 1.0 + i * 10.8, 2.55, 9.5, 6.7, h, items, c, WHITE)
    card(s, 3.6, 11.0, 26.3, 2.6, "答辩表述重点", ["系统不是单独的挂号页面，而是“排班号源 → 预约挂号 → 候诊叫号 → 接诊检查 → 处方发药 → 收费统计”的门诊基础业务链。"], NAVY, rgb("EEF6F7"))
    footer(s, 4)


def structure_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "系统功能结构图", "按老师示例重绘为横向主模块 + 竖排子功能的树状结构")
    s.shapes.add_picture(str(OUT_STRUCTURE), Cm(0.75), Cm(2.2), width=Cm(32.3))
    footer(s, 5)


def flow_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "门诊业务主流程", "用一条主线说明系统如何从患者预约走到收费统计")
    steps = [("患者建档", "基础信息"), ("排班号源", "医生/日期/余号"), ("预约挂号", "生成挂号单"), ("候诊叫号", "进入医生队列"), ("接诊检查", "病历/检查单"), ("处方发药", "审核/出库"), ("收费结算", "账单/支付"), ("费用统计", "收入/图表")]
    x0, y0, w, h, gap = 1.0, 4.0, 3.55, 2.28, 0.48
    for i, (h1, h2) in enumerate(steps, 1):
        x = x0 + (i - 1) * (w + gap)
        card(s, x, y0, w, h, h1, [h2], [TEAL, GREEN, AMBER, NAVY][(i - 1) % 4], WHITE, i)
        if i < len(steps):
            line(s, x + w, y0 + h / 2, x + w + gap, y0 + h / 2)
    lanes = [("患者/挂号员", "患者信息、挂号单、预约时段"), ("医生/检查人员", "候诊队列、病历记录、检查结果"), ("药房/收费员", "处方状态、库存出库、账单支付"), ("管理员/科主任", "科室医生、排班规则、统计看板、日志审计")]
    for i, (r, data) in enumerate(lanes):
        y = 9.1 + i * 1.3
        card(s, 2.0, y, 6.8, 0.82, r, None, TEAL, rgb("EEF7F8"), heading_size=10)
        txt(s, data, 9.4, y + 0.08, 22.0, 0.58, 10.8, INK)
    footer(s, 6)


def architecture_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "技术架构：Qt/C++ 客户端 + Qt TCP 服务端 + MySQL", "客户端负责交互，服务端负责权限和业务规则，数据库负责持久化")
    layers = [("客户端表现层", ["Qt Widgets", "患者预约入口", "角色菜单和业务页面"], TEAL), ("网络通信层", ["QTcpSocket", "TCP + JSON Lines", "统一请求/响应"], GREEN), ("服务端业务层", ["HospitalServer", "RequestRouter 鉴权", "业务 Service"], AMBER), ("数据持久层", ["MySQL", "QODBC/QMYSQL", "演示模式兜底"], NAVY)]
    x = 1.0
    for i, (h, items, c) in enumerate(layers):
        card(s, x, 3.0, 7.1, 6.2, h, items, c, WHITE, i + 1)
        if i < len(layers) - 1:
            line(s, x + 7.1, 6.1, x + 8.1, 6.1)
        x += 8.1
    txt(s, "ApiClient → Protocol → HospitalServer → RequestRouter → Service → DatabaseManager → MySQL", 2.0, 11.0, 29.4, 0.82, 14.5, NAVY, True, PP_ALIGN.CENTER, "Consolas")
    card(s, 4.0, 12.4, 25.8, 2.05, "架构优势", ["客户端不直接操作数据库，内部权限由服务端统一控制；业务数据集中存储，便于统计、审计和后续维护。"], TEAL, rgb("EEF7F8"))
    footer(s, 7)


def database_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "数据库设计", "围绕角色、基础资料、诊疗业务、药品收费四组核心数据建模")
    groups = [("账号权限", ["roles", "users", "operation_logs", "audit_log_details"], TEAL), ("基础资料", ["departments", "doctors", "patients", "doctor_schedules"], GREEN), ("诊疗业务", ["registrations", "medical_records", "examinations", "prescriptions"], AMBER), ("药品收费", ["drug_categories", "drugs", "prescription_items", "stock_records", "bills", "payments", "fee_statistics_daily"], NAVY)]
    x = 1.0
    for i, (h, items, c) in enumerate(groups):
        card(s, x, 2.7, 7.4, 6.75, h, items, c, WHITE)
        if i < len(groups) - 1:
            line(s, x + 7.4, 5.95, x + 8.2, 5.95)
        x += 8.2
    notes = [("患者/医生/排班\n生成挂号记录", TEAL), ("挂号关联\n病历、检查、处方", GREEN), ("处方关联库存\n账单支付完成闭环", AMBER)]
    for i, (n, c) in enumerate(notes, 1):
        card(s, 2.0 + (i - 1) * 10.5, 11.0, 9.2, 1.7, n, None, c, rgb("F8FBFD"), i, heading_size=10.5)
    footer(s, 8)


def implementation_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "核心实现方案", "用项目中的真实工程结构支撑答辩")
    items = [("通信协议", ["QTcpSocket 长连接", "JSON Lines 请求/响应", "module + action 定位接口"], TEAL), ("权限路由", ["RequestRouter 统一鉴权", "按角色控制模块动作", "公开接口限制在预约相关功能"], GREEN), ("业务服务", ["Patient/Registration/Schedule", "Consultation/Examination", "Prescription/Inventory/Billing"], AMBER), ("数据访问", ["DatabaseManager 连接 MySQL", "QODBC/QMYSQL", "连接异常提示与演示模式"], NAVY)]
    for i, (h, body, c) in enumerate(items, 1):
        card(s, 1.0 + (i - 1) * 8.1, 2.7, 7.1, 5.4, h, body, c, WHITE, i)
    card(s, 3.8, 10.5, 26.0, 2.3, "答辩强调点", ["患者预约、医生排班、药品入库等核心写入均可落到 MySQL；页面支持查询、分页、刷新和状态维护，不是固定演示数据。"], TEAL, rgb("EEF7F8"))
    footer(s, 9)


def demo_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "答辩演示路线", "现场演示时用一条最稳的主流程体现系统完整度")
    steps = [("1", "患者预约", "选择科室、医生和时段"), ("2", "挂号管理", "医院端刷新查看挂号单"), ("3", "候诊叫号", "更新患者候诊状态"), ("4", "医生接诊", "填写病历和诊断"), ("5", "处方药房", "审核处方并发药"), ("6", "收费统计", "账单支付和收入统计")]
    for i, (n, h, b) in enumerate(steps):
        x = 1.0 + (i % 3) * 10.8
        y = 3.0 + (i // 3) * 4.2
        card(s, x, y, 9.4, 3.0, h, [b], [TEAL, GREEN, AMBER, NAVY, TEAL, GREEN][i], WHITE, n)
    card(s, 3.5, 12.0, 26.3, 2.2, "演示原则", ["优先展示已经稳定的数据链路；如果现场 MySQL 环境异常，可说明系统保留演示模式，用于展示界面和流程。"], NAVY, rgb("EEF7F8"))
    footer(s, 10)


def plan_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "工作方案与进度安排", "按软件工程过程推进，每个阶段都有明确产出")
    periods = [("需求与开题", "2025.12", "调研/需求/模块"), ("总体设计", "2026.01", "架构/数据库/原型"), ("核心开发", "2026.02-03", "挂号/排班/接诊"), ("扩展开发", "2026.03-04", "检查/处方/收费"), ("测试部署", "2026.04-05", "联调/多环境验证"), ("论文答辩", "2026.05-06", "论文/PPT/提交")]
    for i, (h, t, d) in enumerate(periods, 1):
        x = 1.0 + (i - 1) * 5.35
        card(s, x, 4.0, 4.65, 4.45, h, [t, d], [TEAL, GREEN, AMBER, NAVY, RED, TEAL][i - 1], WHITE, i)
        if i < len(periods):
            line(s, x + 4.65, 6.25, x + 5.25, 6.25)
    card(s, 3.4, 11.0, 27.0, 2.55, "测试重点", ["按患者、挂号员、医生、药房人员、收费员、科主任、管理员等角色分别验证；重点检查号源扣减、候诊状态、处方发药、库存变化、收费统计和操作日志。"], GREEN, rgb("EEF7F8"))
    footer(s, 11)


def outcomes_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "预期成果与项目特色", "答辩时突出完整业务链和工程化实现")
    card(s, 1.0, 2.65, 9.6, 6.6, "预期成果", ["可运行的 C/S 门诊管理系统", "MySQL 建库脚本和演示数据", "跨平台构建与运行脚本", "开题报告、论文和答辩材料"], TEAL)
    card(s, 12.0, 2.65, 9.6, 6.6, "项目特色", ["患者入口与医院人员入口分离", "角色权限驱动菜单和接口", "排班号源影响预约挂号", "处方、药品、收费形成联动"], GREEN)
    card(s, 23.0, 2.65, 9.0, 6.6, "可展示亮点", ["候诊队列与叫号状态", "扫码入库与库存预警", "费用统计和图表展示", "操作日志与审计追踪"], AMBER)
    footer(s, 12)


def risk_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "风险分析与解决措施", "提前准备答辩老师可能关注的问题")
    risks = [("数据库环境差异", ["配置文件管理连接参数", "支持 QODBC/QMYSQL 两种方向", "保留演示模式验证流程"], RED), ("业务状态不同步", ["挂号成功扣减号源", "列表支持手动刷新和自动刷新", "核心状态由服务端统一维护"], AMBER), ("模块多导致测试遗漏", ["按角色设计测试用例", "从建档到收费做端到端验证", "操作日志辅助定位问题"], TEAL)]
    for i, (h, body, c) in enumerate(risks, 1):
        card(s, 1.2 + (i - 1) * 10.9, 3.0, 9.4, 7.2, h, body, c, WHITE, i)
    card(s, 4.0, 12.1, 25.8, 2.2, "应对策略", ["系统功能优先保证主流程稳定，再补充统计、日志、自动刷新等增强能力；答辩演示以稳定链路为主。"], NAVY, rgb("EEF7F8"))
    footer(s, 13)


def close_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, True)
    txt(s, "请各位老师批评指正", 5.0, 6.25, 23.8, 1.3, 34, WHITE, True, PP_ALIGN.CENTER)
    txt(s, TITLE, 5.0, 8.15, 23.8, 0.8, 14, rgb("C9DDE8"), False, PP_ALIGN.CENTER)
    txt(s, STUDENT, 5.0, 9.4, 23.8, 0.65, 11, rgb("DCE9EF"), False, PP_ALIGN.CENTER)
    footer(s, 14)


def build_ppt():
    draw_tree_structure(OUT_STRUCTURE)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    background_slide(prs)
    goal_slide(prs)
    scope_slide(prs)
    structure_slide(prs)
    flow_slide(prs)
    architecture_slide(prs)
    database_slide(prs)
    implementation_slide(prs)
    demo_slide(prs)
    plan_slide(prs)
    outcomes_slide(prs)
    risk_slide(prs)
    close_slide(prs)
    prs.save(str(OUT_PPT))
    shutil.copy2(OUT_PPT, OUT_COPY)


def build_notes():
    OUT_NOTES.write_text(
        """# 开题答辩讲稿提纲（高质量版）

1. 开场说明课题：基于 Qt 和 MySQL 的医院门诊挂号与药品管理系统，采用 C/S 架构。
2. 背景强调：门诊业务高频，挂号、候诊、接诊、处方、库存、收费之间存在强关联，需要数字化闭环。
3. 目标说明：患者入口和医院人员入口分离，内部人员按角色使用不同模块。
4. 功能结构：按 12 个主模块、36 个子功能讲，重点突出系统功能结构图与老师示例一致。
5. 业务流程：排班号源 → 预约挂号 → 候诊叫号 → 医生接诊 → 检查/处方 → 药房发药 → 收费统计。
6. 技术架构：Qt Widgets 客户端负责交互；Qt TCP 服务端负责权限和业务规则；MySQL 负责业务数据持久化。
7. 数据库设计：角色、用户、科室、医生、患者、排班、挂号、病历、检查、处方、药品、库存、账单、支付、统计、日志。
8. 演示路线：患者预约、医院端刷新挂号、候诊叫号、医生接诊、处方发药、收费统计。
9. 风险回答：数据库环境、业务状态同步、模块测试覆盖。
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    build_ppt()
    build_notes()
    print(OUT_PPT)
    print(OUT_COPY)
    print(OUT_STRUCTURE)
    print(OUT_NOTES)

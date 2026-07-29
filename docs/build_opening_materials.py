from pathlib import Path
import shutil

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm as PCm, Inches, Pt as PPt


ROOT = Path(__file__).resolve().parent
PROJECT_ROOT = ROOT.parent
DESKTOP_REPORT = Path("E:/15751/Desktop/开题报告.docx")
OUT_REPORT = Path("E:/15751/Desktop/开题报告-补充完善版.docx")
OUT_REPORT_DOCS = ROOT / "开题报告-补充完善版.docx"
OUT_PPT = Path("E:/15751/Desktop/医院门诊挂号与药品管理系统-开题答辩PPT.pptx")
OUT_PPT_DOCS = ROOT / "医院门诊挂号与药品管理系统-开题答辩PPT.pptx"
OUT_STRUCTURE = ROOT / "系统功能结构图-答辩版.png"
OUT_ARCH = ROOT / "系统技术架构图-答辩版.png"
OUT_FLOW = ROOT / "门诊业务流程图-答辩版.png"
OUT_DB = ROOT / "数据库核心关系图-答辩版.png"
OUT_MMD = ROOT / "系统功能结构图-答辩版.mmd"
OUT_NOTES = ROOT / "开题答辩讲稿提纲.md"

TITLE = "基于 Qt 和 MySQL 的医院门诊挂号与药品管理系统"
SUBTITLE = "开题答辩"

BLUE = RGBColor(32, 88, 167)
NAVY = RGBColor(17, 34, 64)
TEAL = RGBColor(28, 132, 128)
GREEN = RGBColor(59, 130, 81)
ORANGE = RGBColor(206, 110, 42)
RED = RGBColor(180, 65, 61)
GRAY = RGBColor(90, 101, 116)
LIGHT_BG = RGBColor(244, 248, 252)
WHITE = RGBColor(255, 255, 255)


def font(size, bold=False):
    candidates = [
        Path(r"C:\Windows\Fonts\simhei.ttf") if bold else Path(r"C:\Windows\Fonts\simsun.ttc"),
        Path(r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
    ]
    for item in candidates:
        if item.exists():
            return ImageFont.truetype(str(item), size)
    return ImageFont.load_default()


def wrapped_lines(draw, text, font_obj, max_width):
    lines = []
    for paragraph in text.split("\n"):
        current = ""
        for ch in paragraph:
            test = current + ch
            bbox = draw.textbbox((0, 0), test, font=font_obj)
            if bbox[2] - bbox[0] <= max_width or not current:
                current = test
            else:
                lines.append(current)
                current = ch
        lines.append(current)
    return lines


def draw_text_box(draw, box, text, font_obj, fill=(20, 32, 46), align="center", spacing=8):
    x1, y1, x2, y2 = box
    lines = wrapped_lines(draw, text, font_obj, max(10, x2 - x1 - 26))
    line_heights = []
    widths = []
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font_obj)
        widths.append(bbox[2] - bbox[0])
        line_heights.append(bbox[3] - bbox[1] + spacing)
    total_h = sum(line_heights) - spacing if line_heights else 0
    y = y1 + max(0, (y2 - y1 - total_h) / 2)
    for i, line in enumerate(lines):
        w = widths[i]
        if align == "left":
            x = x1 + 18
        else:
            x = x1 + (x2 - x1 - w) / 2
        draw.text((x, y), line, font=font_obj, fill=fill)
        y += line_heights[i]


def arrow(draw, start, end, color=(67, 88, 118), width=4):
    draw.line((start, end), fill=color, width=width)
    sx, sy = start
    ex, ey = end
    if abs(ex - sx) >= abs(ey - sy):
        direction = 1 if ex > sx else -1
        points = [(ex, ey), (ex - 16 * direction, ey - 9), (ex - 16 * direction, ey + 9)]
    else:
        direction = 1 if ey > sy else -1
        points = [(ex, ey), (ex - 9, ey - 16 * direction), (ex + 9, ey - 16 * direction)]
    draw.polygon(points, fill=color)


def draw_structure(path):
    w, h = 3200, 1500
    img = Image.new("RGB", (w, h), "white")
    draw = ImageDraw.Draw(img)
    title_font = font(54, True)
    module_font = font(32, True)
    leaf_font = font(27)
    caption_font = font(30)
    line = (63, 84, 112)
    root_fill = (226, 238, 249)
    module_fill = (234, 244, 241)
    leaf_fill = (250, 252, 255)
    border = (54, 74, 102)

    root = (910, 55, 2290, 135)
    draw.rounded_rectangle(root, radius=12, fill=root_fill, outline=border, width=3)
    draw_text_box(draw, root, "医院门诊挂号与药品管理系统", title_font, fill=(15, 34, 60))

    modules = [
        ("登录权限", ["患者入口", "人员登录", "角色权限", "操作日志"]),
        ("院长驾驶舱", ["汇总指标", "待办概览", "收入看板"]),
        ("患者管理", ["患者建档", "信息查询", "病历档案"]),
        ("科室医生", ["科室维护", "医生信息", "在职状态"]),
        ("排班号源", ["排班维护", "号源调整", "停诊管理"]),
        ("挂号候诊", ["预约挂号", "挂号管理", "候诊叫号"]),
        ("诊疗检查", ["医生接诊", "病历记录", "检查检验"]),
        ("处方药房", ["处方开立", "处方审核", "确认发药"]),
        ("药品库存", ["药品维护", "扫码入库", "库存预警"]),
        ("收费结算", ["账单查询", "缴费退费", "支付记录"]),
        ("统计分析", ["日收入", "科室统计", "图表展示"]),
        ("系统运维", ["配置管理", "MySQL连接", "多端运行"]),
    ]
    cols = 6
    card_w, card_h = 470, 385
    gap_x, gap_y = 42, 100
    start_x, start_y = 75, 300
    trunk_y = 215
    draw.line((w // 2, root[3], w // 2, trunk_y), fill=line, width=4)
    draw.line((start_x + card_w // 2, trunk_y, start_x + 5 * (card_w + gap_x) + card_w // 2, trunk_y), fill=line, width=4)

    for idx, (name, leaves) in enumerate(modules):
        row, col = divmod(idx, cols)
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        cx = x + card_w // 2
        draw.line((cx, trunk_y, cx, y), fill=line, width=4)
        card = (x, y, x + card_w, y + card_h)
        draw.rounded_rectangle(card, radius=10, fill=module_fill, outline=border, width=3)
        head = (x, y, x + card_w, y + 75)
        draw.rounded_rectangle(head, radius=10, fill=(213, 232, 228), outline=border, width=0)
        draw_text_box(draw, head, name, module_font, fill=(20, 64, 66))
        for j, item in enumerate(leaves):
            leaf = (x + 42, y + 108 + j * 66, x + card_w - 42, y + 158 + j * 66)
            draw.rounded_rectangle(leaf, radius=6, fill=leaf_fill, outline=(169, 185, 202), width=2)
            draw_text_box(draw, leaf, item, leaf_font, fill=(30, 43, 60))

    draw.text((w // 2 - 145, 1375), "图1  系统功能结构图", font=caption_font, fill=(20, 32, 46))
    img.save(path)


def draw_architecture(path):
    w, h = 2200, 1150
    img = Image.new("RGB", (w, h), "white")
    draw = ImageDraw.Draw(img)
    title_font = font(46, True)
    box_font = font(30, True)
    item_font = font(25)
    colors = [(232, 240, 253), (232, 246, 242), (255, 244, 232), (245, 235, 250)]
    border = (52, 73, 102)
    draw.text((70, 55), "系统技术架构", font=title_font, fill=(17, 34, 64))

    boxes = [
        ((95, 190, 550, 880), "客户端表现层", ["Qt Widgets 桌面界面", "患者预约挂号入口", "医院人员角色菜单", "表格分页 / 查询 / 刷新"], colors[0]),
        ((690, 190, 1145, 880), "网络通信层", ["QTcpSocket", "TCP + JSON Lines", "请求模块/动作路由", "统一响应与错误提示"], colors[1]),
        ((1285, 190, 1740, 880), "服务端业务层", ["RequestRouter 权限校验", "Auth/Patient/Registration", "Schedule/Consultation", "Prescription/Billing/Stats"], colors[2]),
        ((1885, 190, 2160, 880), "数据持久层", ["MySQL", "20 张左右业务表", "QODBC/QMYSQL", "演示模式兜底"], colors[3]),
    ]
    for box, title, items, fill in boxes:
        draw.rounded_rectangle(box, radius=18, fill=fill, outline=border, width=3)
        draw_text_box(draw, (box[0] + 20, box[1] + 28, box[2] - 20, box[1] + 92), title, box_font, fill=(15, 34, 60))
        for i, item in enumerate(items):
            y = box[1] + 145 + i * 95
            draw.rounded_rectangle((box[0] + 45, y, box[2] - 45, y + 58), radius=8, fill=(255, 255, 255), outline=(177, 190, 208), width=2)
            draw_text_box(draw, (box[0] + 55, y, box[2] - 55, y + 58), item, item_font, fill=(34, 49, 68))
    for x in [550, 1145, 1740]:
        arrow(draw, (x + 25, 535), (x + 135, 535))
    draw_text_box(draw, (385, 970, 1815, 1060), "设计要点：客户端只负责交互，服务端统一处理权限与业务规则，数据库集中保存门诊业务数据。", item_font, fill=(60, 72, 88))
    img.save(path)


def draw_flow(path):
    w, h = 2200, 1000
    img = Image.new("RGB", (w, h), "white")
    draw = ImageDraw.Draw(img)
    title_font = font(44, True)
    box_font = font(27, True)
    small_font = font(22)
    draw.text((70, 48), "门诊业务主流程", font=title_font, fill=(17, 34, 64))
    steps = [
        ("患者建档", "维护患者基础信息"),
        ("排班号源", "医生/日期/时段/余号"),
        ("预约挂号", "生成挂号单并扣减号源"),
        ("候诊叫号", "按状态进入医生队列"),
        ("接诊检查", "病历、诊断、检查单"),
        ("处方发药", "处方审核与库存出库"),
        ("收费结算", "账单、支付、退费"),
        ("统计日志", "收入统计和操作追溯"),
    ]
    x0, y0 = 70, 210
    step_w, step_h, gap = 235, 150, 35
    for i, (name, desc) in enumerate(steps):
        x = x0 + i * (step_w + gap)
        color = [(229, 240, 255), (233, 246, 241), (255, 244, 231), (244, 237, 250)][i % 4]
        draw.rounded_rectangle((x, y0, x + step_w, y0 + step_h), radius=14, fill=color, outline=(66, 86, 116), width=3)
        draw_text_box(draw, (x + 12, y0 + 22, x + step_w - 12, y0 + 72), name, box_font, fill=(19, 42, 70))
        draw_text_box(draw, (x + 16, y0 + 84, x + step_w - 16, y0 + 132), desc, small_font, fill=(55, 66, 82))
        if i < len(steps) - 1:
            arrow(draw, (x + step_w + 5, y0 + step_h // 2), (x + step_w + gap - 8, y0 + step_h // 2), width=3)

    lanes = [
        ("患者/挂号员", "患者信息、挂号单、预约时段"),
        ("医生/检查人员", "候诊队列、病历记录、检查结果"),
        ("药房/收费员", "处方状态、药品库存、账单支付"),
        ("管理员/科主任", "科室医生、排班规则、统计看板、日志审计"),
    ]
    for i, (role, data) in enumerate(lanes):
        y = 520 + i * 95
        draw.rounded_rectangle((120, y, 2020, y + 62), radius=10, fill=(248, 250, 252), outline=(180, 193, 210), width=2)
        draw_text_box(draw, (145, y, 410, y + 62), role, box_font, fill=(28, 91, 146))
        draw_text_box(draw, (445, y, 1960, y + 62), data, small_font, fill=(53, 65, 82), align="left")
    img.save(path)


def draw_database(path):
    w, h = 2200, 1180
    img = Image.new("RGB", (w, h), "white")
    draw = ImageDraw.Draw(img)
    title_font = font(44, True)
    table_font = font(25, True)
    item_font = font(20)
    draw.text((70, 45), "数据库核心关系", font=title_font, fill=(17, 34, 64))

    groups = [
        ((90, 170, 500, 450), "账号权限", ["roles", "users", "operation_logs", "audit_log_details"], (232, 240, 253)),
        ((620, 170, 1030, 450), "基础资料", ["departments", "doctors", "patients", "doctor_schedules"], (233, 246, 241)),
        ((1150, 170, 1560, 450), "诊疗业务", ["registrations", "medical_records", "examinations", "prescriptions"], (255, 244, 232)),
        ((1680, 170, 2110, 450), "药品收费", ["drug_categories", "drugs", "prescription_items", "stock_records", "bills", "payments", "fee_statistics_daily"], (245, 235, 250)),
    ]
    for box, title, tables, color in groups:
        draw.rounded_rectangle(box, radius=14, fill=color, outline=(65, 84, 112), width=3)
        draw_text_box(draw, (box[0] + 12, box[1] + 22, box[2] - 12, box[1] + 70), title, table_font, fill=(20, 42, 70))
        for i, name in enumerate(tables):
            y = box[1] + 88 + i * 35
            draw.text((box[0] + 34, y), f"• {name}", font=item_font, fill=(43, 55, 74))

    relations = [
        ((500, 310), (620, 310), "用户关联医生/操作人"),
        ((1030, 310), (1150, 310), "排班驱动挂号"),
        ((1560, 310), (1680, 310), "处方关联药品/账单"),
        ((1810, 450), (1810, 650), "收入汇总"),
        ((290, 450), (290, 650), "审计追踪"),
    ]
    for s, e, label in relations:
        arrow(draw, s, e, width=3)
        mx, my = (s[0] + e[0]) // 2, (s[1] + e[1]) // 2
        draw.text((mx - 95, my - 35), label, font=item_font, fill=(73, 86, 105))

    notes = [
        "关键约束：挂号记录 registration_no、处方 prescription_no、账单 bill_no 等业务编号唯一。",
        "主外键链路：患者/医生/排班 → 挂号 → 病历/检查/处方 → 药品库存/收费结算。",
        "统计与日志：费用统计按日期和科室沉淀，操作日志记录核心写操作，便于追溯。",
    ]
    for i, note in enumerate(notes):
        y = 725 + i * 82
        draw.rounded_rectangle((160, y, 2040, y + 56), radius=8, fill=(248, 250, 252), outline=(185, 196, 211), width=2)
        draw_text_box(draw, (185, y, 2000, y + 56), note, item_font, fill=(42, 54, 72), align="left")
    img.save(path)


def write_mermaid(path):
    path.write_text(
        """graph TD
    A[医院门诊挂号与药品管理系统]
    A --> B[登录权限]
    A --> C[院长驾驶舱]
    A --> D[患者管理]
    A --> E[科室医生]
    A --> F[排班号源]
    A --> G[挂号候诊]
    A --> H[诊疗检查]
    A --> I[处方药房]
    A --> J[药品库存]
    A --> K[收费结算]
    A --> L[统计分析]
    A --> M[系统运维]
    B --> B1[患者入口]
    B --> B2[人员登录]
    B --> B3[角色权限]
    G --> G1[预约挂号]
    G --> G2[挂号管理]
    G --> G3[候诊叫号]
    H --> H1[医生接诊]
    H --> H2[病历记录]
    H --> H3[检查检验]
    I --> I1[处方开立]
    I --> I2[处方审核]
    I --> I3[确认发药]
    J --> J1[扫码入库]
    J --> J2[库存预警]
    K --> K1[账单查询]
    K --> K2[缴费退费]
    L --> L1[日收入]
    L --> L2[科室统计]
""",
        encoding="utf-8",
    )


BACKGROUND = """研究背景与意义：
（一）研究背景
随着公立医院高质量发展和智慧医院建设持续推进，门诊服务正在从传统窗口办理逐步转向线上线下一体化协同。国务院办公厅在《关于推动公立医院高质量发展的意见》中提出，要推动云计算、大数据、物联网等新一代信息技术与医疗服务深度融合；国家卫生健康委等部门发布的《公立医院高质量发展促进行动（2021—2025年）》也将信息化建设作为提升医疗服务和医院管理能力的重要支撑[1][2]。在门诊场景中，患者就诊通常经过患者建档、科室与医生选择、预约挂号、候诊叫号、医生接诊、检查检验、处方开立、药品发放、收费结算和费用统计等环节。各环节数据关联紧密、状态变化频繁，如果仍依赖纸质单据或简单表格管理，容易出现重复录入、号源不同步、库存更新滞后、收费统计不准确和责任追溯困难等问题。
从业务需求看，挂号系统不应只停留在“保存一条预约记录”，还需要与医生排班、时段号源、候诊队列和收费账单保持一致；药品管理也不只是登记库存数量，还应覆盖药品基础信息、分类、扫码入库、库存预警、处方审核与出库记录。门诊业务的典型特点是流程短、频次高、窗口并发多、数据变更快，因此系统设计需要在可用性、准确性和可维护性之间取得平衡。本课题选择医院门诊挂号与药品管理作为研究对象，能够覆盖医院信息系统中较典型的患者服务、诊疗协同、药品流转和费用管理流程。
从技术实现看，Qt 具备成熟的跨平台桌面界面开发能力，Qt Widgets 适合构建稳定、响应清晰的业务型桌面应用；Qt SQL 可通过 QODBC 或 QMYSQL 访问 MySQL 数据库，便于将界面操作、业务规则和数据持久化连接起来[5][6]。MySQL 作为成熟的关系型数据库，适合保存患者、医生、排班、挂号、病历、处方、药品、库存、账单、支付、统计和日志等结构化数据[7]。本项目采用 Qt/C++ 实现客户端和服务端，服务端通过 TCP + JSON Lines 协议提供业务接口和权限校验，数据库负责持久化保存业务数据，能够体现客户端、服务端、数据库和工程构建的综合实践能力。
（二）研究意义
本课题的意义主要体现在三个方面。第一，在业务应用层面，系统围绕门诊主流程设计，将患者信息、医生排班、挂号号源、候诊叫号、接诊病历、检查检验、处方药品、库存变化、收费结算和统计分析统一管理，有助于减少信息孤岛和重复录入，提高门诊窗口处理效率。第二，在数据库设计层面，系统通过角色、用户、科室、医生、患者、排班、挂号、病历、检查、处方、药品、库存、账单、支付、统计和审计日志等表建立主外键关系，能够训练较完整的关系数据库建模、查询、事务和数据一致性设计能力。第三，在工程实践层面，本课题需要完成 Qt 界面设计、网络通信、服务端路由、角色权限控制、MySQL 访问、分页查询、自动刷新、统计图表、异常处理和跨平台构建，能够完整体现软件工程中的需求分析、概要设计、详细设计、编码实现、测试验证和部署维护过程。与单一模块的小型管理程序相比，本系统功能链条更完整，更接近真实医院门诊信息系统的基础形态，具有较好的实践价值和教学价值。
"""

MAIN_CONTENT = """主要内容：
本系统面向医院门诊挂号与药品管理业务，采用客户端/服务端结构。客户端负责患者预约入口和医院人员业务界面，服务端负责统一接口、权限校验和业务处理，MySQL 数据库负责保存核心业务数据。系统整体功能结构图如图1所示。

一、登录权限与入口管理模块
系统区分“患者预约挂号”和“医院人员登录”两类入口。患者预约无需医院内部账号，可查询公开排班并提交预约；医院人员使用账号密码登录，按照管理员、科主任、挂号员、医生、药房人员、收费员等角色展示不同菜单。服务端 RequestRouter 对模块和动作进行权限校验，并对新增、修改、删除、叫号、审核、收费等关键操作记录日志。
二、院长驾驶舱与统计概览模块
该模块为管理员、科主任和收费人员提供门诊核心指标汇总，包括今日挂号、待缴费、已缴费、药品库存预警和收入统计等内容，便于管理者快速了解门诊运行状态。
三、患者管理与病历档案模块
患者管理维护患者编号、姓名、性别、出生日期、身份证号、联系电话、地址等基础资料，支持新增、修改、删除、查询和分页展示。患者病历档案关联挂号、接诊、检查和处方记录，便于医生回看患者历史就诊情况。
四、科室与医生管理模块
科室管理维护门诊科室、专科和诊室目录；医生管理维护医生所属科室、职称、专长、挂号费和在职状态。该模块为排班、挂号和统计提供基础数据，保证患者端看到的科室医生信息来自真实数据库。
五、医生排班与号源管理模块
医生排班按医生、日期和时段维护出诊安排，设置总号源、剩余号源和出诊状态。排班数据直接影响患者端可预约医生列表，挂号成功后扣减对应号源，避免固定演示数据造成业务失真。
六、挂号管理与候诊叫号模块
挂号管理完成患者、科室、医生、日期和时段选择，生成挂号记录并维护等待、已叫号、检查中、已完成、已取消等状态。候诊队列支持按科室、医生和状态筛选，医生或挂号员可进行叫号操作，增强门诊流程的连续性。
七、医生接诊、检查检验与处方模块
医生接诊模块展示待诊患者，医生可填写主诉、现病史、既往史、体征、诊断和医嘱。检查检验模块支持开立检查单、录入检查结果和查看状态。处方管理支持开立处方、添加药品明细、计算金额、处方审核和确认发药，与药品库存和收费模块形成业务联动。
八、药品库存管理模块
药品库存维护药品编码、条形码、名称、分类、规格、单位、进价、售价、库存数量、预警数量和有效期，支持药品查询、扫码入库、手动入库、库存预警和出入库记录追踪。对于已有药品，入库时自动增加库存；对于新药品，可补充分类和基础信息。
九、收费结算模块
收费结算根据挂号费、药品费和检查费生成账单，支持账单查询、缴费、退费、支付方式记录和收费状态维护。模块与挂号、处方和检查数据关联，减少人工统计错误。
十、系统运维与多端运行模块
项目使用 CMake 组织 client、server、common、launcher 等子工程，支持 VS Code、Qt Creator、Windows、Linux 和银河麒麟环境构建运行。服务端通过配置文件选择 MySQL 连接参数，也保留演示模式，便于课堂演示和环境异常时验证系统功能。
"""

PLAN = """工作方案及进度安排：
（一）工作方案
本课题按照软件工程开发流程推进。首先进行需求分析，梳理患者预约、医院人员登录、医生排班、挂号候诊、医生接诊、检查检验、处方药房、药品库存、收费结算、统计分析和操作日志等功能边界；其次进行数据库设计，建立角色、用户、科室、医生、患者、排班、挂号、病历、检查、处方、药品、库存、账单、支付、统计和审计日志等数据表，并确定主外键关系与核心业务编号；然后进行系统架构设计，采用 Qt Widgets 实现客户端界面，采用 Qt TCP 服务端提供接口，使用 JSON Lines 作为网络消息格式，使用 MySQL 进行数据持久化；最后进行编码实现、联调测试、部署验证和论文撰写。
技术路线为：Qt Widgets 客户端界面 → ApiClient 封装 TCP 请求 → Protocol 编解码 JSON Lines → HospitalServer 接收连接 → RequestRouter 进行权限校验和服务分发 → 各业务 Service 处理请求 → DatabaseManager 连接 MySQL → 返回统一响应。测试阶段将按照患者、挂号员、医生、药房人员、收费员、科主任和管理员等角色分别验证，重点关注号源扣减是否正确、挂号新增后医院端是否刷新显示、处方与药品库存是否联动、账单支付状态是否一致、统计数据是否可追溯、断开数据库后的错误提示是否清晰。
（二）进度安排
2025年12月08日—2025年12月21日：完成课题调研、需求分析、系统功能模块划分和开题报告撰写。
2025年12月22日—2026年01月04日：完成系统总体设计，包括 C/S 架构、数据库表结构、主要业务流程和界面原型。
2026年01月05日—2026年02月08日：完成项目基础框架，实现客户端登录入口、服务端通信框架、MySQL 连接配置和基础数据初始化。
2026年02月09日—2026年03月01日：完成患者管理、挂号预约、挂号管理、候诊队列和医生排班等核心模块，实现号源限制、刷新和分页查询。
2026年03月02日—2026年03月22日：完成科室医生管理、医生接诊、病历档案、检查检验和处方管理模块，实现诊疗业务闭环。
2026年03月23日—2026年04月12日：完成药品库存、扫码入库、处方审核发药、收费结算和费用统计模块，实现库存、账单和统计联动。
2026年04月13日—2026年04月26日：进行系统联调和功能测试，修复权限控制、数据库连接、页面刷新、分页显示和异常提示问题。
2026年04月27日—2026年05月10日：完成 Windows、Linux、银河麒麟、Qt Creator 和 VS Code 环境下的构建运行验证，整理部署说明。
2026年05月11日—2026年05月24日：撰写毕业设计论文，完善系统截图、数据库设计说明、核心代码说明和测试分析。
2026年05月25日—2026年06月05日：根据指导教师意见修改论文和系统，准备答辩 PPT，完成最终提交。
"""

REFERENCES = """四、参考文献：
[1] 国务院办公厅. 关于推动公立医院高质量发展的意见[Z]. 2021.
[2] 国家卫生健康委, 国家中医药管理局. 公立医院高质量发展促进行动（2021—2025年）[Z]. 2021.
[3] 国家卫生健康委, 国家中医药管理局, 国家疾病预防控制局. 医疗卫生机构网络安全管理办法[Z]. 2022.
[4] 国家卫生健康委. 三级医院评审标准（2022年版）及其实施细则[S]. 2022.
[5] The Qt Company. Qt Widgets Documentation[EB/OL]. https://doc.qt.io/qt-6/qtwidgets-index.html, 2026.
[6] The Qt Company. QSqlDatabase Class Documentation[EB/OL]. https://doc.qt.io/qt-6/qsqldatabase.html, 2026.
[7] Oracle Corporation. MySQL 8.4 Reference Manual[EB/OL]. https://dev.mysql.com/doc/refman/8.4/en/, 2026.
[8] Kitware. CMake Documentation[EB/OL]. https://cmake.org/documentation/, 2026.
[9] HL7 International. FHIR Release 5 Specification[EB/OL]. https://hl7.org/fhir/, 2023.
[10] 王珊, 萨师煊. 数据库系统概论[M]. 第6版. 北京: 高等教育出版社, 2023.
[11] 张海藩, 牟永敏. 软件工程导论[M]. 第7版. 北京: 清华大学出版社, 2022.
"""


def set_cell_text(cell, text, bold_first=True, font_size=10.5):
    cell.text = ""
    for idx, part in enumerate(text.split("\n")):
        p = cell.paragraphs[0] if idx == 0 else cell.add_paragraph()
        p.paragraph_format.line_spacing = 1.22
        p.paragraph_format.space_after = Pt(2)
        stripped = part.strip()
        no_indent_prefixes = ("一、", "二、", "三、", "四、", "五、", "六、", "七、", "八、", "九、", "十、", "（", "[", "202", "图")
        p.paragraph_format.first_line_indent = Pt(21) if stripped and not stripped.startswith(no_indent_prefixes) else Pt(0)
        run = p.add_run(part)
        run.font.name = "宋体"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
        run.font.size = Pt(font_size)
        if bold_first and idx == 0:
            run.bold = True


def set_short_cell(cell, text):
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(text)
    run.font.name = "宋体"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
    run.font.size = Pt(10.5)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def set_table_borders(table):
    borders = table._tbl.tblPr.first_child_found_in("w:tblBorders")
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        table._tbl.tblPr.append(borders)
    for name in ("top", "left", "bottom", "right", "insideH", "insideV"):
        elem = borders.find(qn(f"w:{name}"))
        if elem is None:
            elem = OxmlElement(f"w:{name}")
            borders.append(elem)
        elem.set(qn("w:val"), "single")
        elem.set(qn("w:sz"), "8")
        elem.set(qn("w:space"), "0")
        elem.set(qn("w:color"), "000000")


def insert_picture(cell, image_path, width_cm=15.8):
    p = cell.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    run.add_picture(str(image_path), width=Cm(width_cm))
    cap = cell.add_paragraph("图1  系统功能结构图")
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in cap.runs:
        run.font.name = "宋体"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
        run.font.size = Pt(10)


def build_report():
    if not DESKTOP_REPORT.exists():
        raise FileNotFoundError(f"未找到原开题报告：{DESKTOP_REPORT}")
    doc = Document(str(DESKTOP_REPORT))
    table = doc.tables[0]
    set_table_borders(table)
    for section in doc.sections:
        section.top_margin = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin = Cm(2.2)
        section.right_margin = Cm(2.2)

    title_cell = table.cell(3, 1).merge(table.cell(3, 3))
    set_short_cell(title_cell, TITLE)

    content_rows = [(4, BACKGROUND), (5, MAIN_CONTENT), (6, PLAN), (7, REFERENCES)]
    for row_idx, content in content_rows:
        cell = table.cell(row_idx, 0).merge(table.cell(row_idx, 3))
        set_cell_text(cell, content)
        if row_idx == 5:
            insert_picture(cell, OUT_STRUCTURE)

    for row_idx in (8, 9):
        cell = table.cell(row_idx, 0).merge(table.cell(row_idx, 3))
        if row_idx == 8:
            set_cell_text(cell, "指导教师意见：\n\n\n\n\n指导教师：\n                                     年  月  日")
        else:
            set_cell_text(cell, "所在专业意见：\n□通过    □不通过\n\n\n负责人：\n                                     年  月  日")

    doc.save(str(OUT_REPORT))
    shutil.copy2(OUT_REPORT, OUT_REPORT_DOCS)


def add_bg(slide, prs):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(246, 249, 252)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PCm(0.28))
    rect.fill.solid()
    rect.fill.fore_color.rgb = BLUE
    rect.line.fill.background()


def add_title(slide, text, subtitle=None):
    title = slide.shapes.add_textbox(PCm(0.9), PCm(0.65), PCm(25.5), PCm(1.0))
    p = title.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "微软雅黑"
    p.font.size = PPt(27)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        st = slide.shapes.add_textbox(PCm(0.95), PCm(1.62), PCm(23), PCm(0.6))
        sp = st.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "微软雅黑"
        sp.font.size = PPt(12)
        sp.font.color.rgb = GRAY


def add_footer(slide, idx):
    tx = slide.shapes.add_textbox(PCm(25.0), PCm(14.4), PCm(1.2), PCm(0.35))
    p = tx.text_frame.paragraphs[0]
    p.text = f"{idx:02d}"
    p.font.name = "Arial"
    p.font.size = PPt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_card(slide, x, y, w, h, title, body, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PCm(x), PCm(y), PCm(w), PCm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 228, 238)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PCm(x), PCm(y), PCm(0.16), PCm(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = PCm(0.35)
    tf.margin_right = PCm(0.25)
    tf.margin_top = PCm(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "微软雅黑"
    p.font.size = PPt(15)
    p.font.bold = True
    p.font.color.rgb = color
    for line in body:
        bp = tf.add_paragraph()
        bp.text = line
        bp.level = 0
        bp.font.name = "微软雅黑"
        bp.font.size = PPt(10.5)
        bp.font.color.rgb = RGBColor(43, 55, 72)
        bp.space_before = PPt(4)
    return shape


def add_bullets(slide, x, y, w, h, bullets, size=14):
    box = slide.shapes.add_textbox(PCm(x), PCm(y), PCm(w), PCm(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "微软雅黑"
        p.font.size = PPt(size)
        p.font.color.rgb = RGBColor(38, 52, 72)
        p.space_after = PPt(6)
    return box


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []
    for _ in range(13):
        slide = prs.slides.add_slide(blank)
        add_bg(slide, prs)
        slides.append(slide)

    # 1
    s = slides[0]
    cover = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    cover.fill.solid()
    cover.fill.fore_color.rgb = RGBColor(239, 245, 250)
    cover.line.fill.background()
    s.shapes.add_picture(str(OUT_FLOW), PCm(6.8), PCm(1.25), width=PCm(19.5))
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PCm(0.0), PCm(0.0), PCm(11.8), PCm(19.1))
    panel.fill.solid()
    panel.fill.fore_color.rgb = NAVY
    panel.line.fill.background()
    title = s.shapes.add_textbox(PCm(1.05), PCm(2.25), PCm(10.2), PCm(2.7))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = TITLE
    p.font.name = "微软雅黑"
    p.font.size = PPt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sub = s.shapes.add_textbox(PCm(1.08), PCm(5.45), PCm(8), PCm(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "开题答辩 | Qt/C++ · MySQL · C/S 架构"
    sp.font.name = "微软雅黑"
    sp.font.size = PPt(14)
    sp.font.color.rgb = RGBColor(197, 215, 239)
    meta = s.shapes.add_textbox(PCm(1.08), PCm(10.5), PCm(8.2), PCm(2.0))
    mp = meta.text_frame.paragraphs[0]
    mp.text = "软件学院  计算机科学与技术\n学生：刘子航    指导教师：李长毅"
    mp.font.name = "微软雅黑"
    mp.font.size = PPt(13)
    mp.font.color.rgb = RGBColor(235, 242, 250)

    # 2
    s = slides[1]
    add_title(s, "课题背景：门诊业务需要数字化闭环", "从窗口型处理转向多角色协同、数据可追溯的门诊业务系统")
    add_card(s, 1.0, 2.4, 7.4, 3.3, "现实问题", ["排队和人工登记导致信息重复录入", "号源、候诊、收费状态变化频繁", "药品库存、处方和账单容易割裂"], BLUE)
    add_card(s, 9.1, 2.4, 7.4, 3.3, "建设方向", ["预约诊疗、智慧医院和精细化管理", "患者服务、医生诊疗、药房收费一体化", "通过日志和统计提升管理可视化"], TEAL)
    add_card(s, 17.2, 2.4, 7.4, 3.3, "课题价值", ["覆盖 Qt 客户端、服务端和数据库", "体现需求分析、设计、实现和测试全过程", "适合毕业设计展示真实工程能力"], GREEN)

    # 3
    s = slides[2]
    add_title(s, "系统目标与使用对象", "目标不是单个挂号窗口，而是门诊挂号、诊疗、药品、收费的基础业务链")
    add_card(s, 1.0, 2.2, 7.2, 4.6, "系统目标", ["患者可预约挂号", "医院人员按角色进入业务模块", "排班号源、挂号候诊、处方库存、收费统计形成联动"], BLUE)
    roles = [("患者", "预约挂号"), ("挂号员", "建档/挂号/叫号"), ("医生", "接诊/病历/检查/处方"), ("药房", "审核/发药/库存"), ("收费员", "账单/支付/统计"), ("管理员", "基础数据/日志")]
    for i, (role, task) in enumerate(roles):
        x = 9.2 + (i % 3) * 5.2
        y = 2.1 + (i // 3) * 2.35
        add_card(s, x, y, 4.55, 1.55, role, [task], [TEAL, BLUE, GREEN, ORANGE, RED, NAVY][i])

    # 4
    s = slides[3]
    add_title(s, "功能边界：15 个医院端页面 + 患者预约入口", "源码已实现清晰的菜单权限和模块划分")
    add_bullets(s, 1.1, 2.25, 8.0, 5.7, [
        "院长驾驶舱、患者管理、挂号管理、候诊队列",
        "患者病历档案、科室管理、医生排班、医生管理",
        "医生接诊、检查检验、处方管理、药品库存",
        "收费结算、费用统计、操作日志"
    ], 14)
    add_card(s, 10.1, 2.0, 6.2, 4.8, "权限设计", ["管理员拥有全部权限", "科主任偏向诊疗和统计", "挂号员负责患者、挂号和排班", "医生负责接诊、检查和处方", "药房负责库存与处方发药", "收费员负责账单和收入统计"], BLUE)
    add_card(s, 17.2, 2.0, 7.1, 4.8, "公开入口", ["患者端可直接查询公开排班", "提交预约后进入医院端挂号列表", "用服务端权限规则保护内部业务"], TEAL)

    # 5
    s = slides[4]
    add_title(s, "系统功能结构图", "根据项目实际页面和服务模块整理")
    s.shapes.add_picture(str(OUT_STRUCTURE), PCm(0.75), PCm(2.05), width=PCm(25.0))

    # 6
    s = slides[5]
    add_title(s, "门诊业务流程", "从患者建档到收费统计，形成可追溯业务链")
    s.shapes.add_picture(str(OUT_FLOW), PCm(0.85), PCm(1.9), width=PCm(24.8))

    # 7
    s = slides[6]
    add_title(s, "技术架构：Qt/C++ 客户端 + Qt TCP 服务端 + MySQL", "客户端专注交互，服务端统一权限和业务规则，数据库集中保存数据")
    s.shapes.add_picture(str(OUT_ARCH), PCm(1.0), PCm(1.85), width=PCm(24.3))

    # 8
    s = slides[7]
    add_title(s, "数据库设计", "围绕角色、基础资料、诊疗业务、药品收费四组核心数据建模")
    s.shapes.add_picture(str(OUT_DB), PCm(1.0), PCm(1.65), width=PCm(24.3))

    # 9
    s = slides[8]
    add_title(s, "核心实现方案", "用项目中的真实工程结构支撑答辩")
    add_card(s, 1.0, 2.05, 5.8, 4.7, "通信协议", ["QTcpSocket 长连接", "JSON Lines 请求/响应", "module + action 定位接口", "统一 success/message/data"], BLUE)
    add_card(s, 7.7, 2.05, 5.8, 4.7, "权限路由", ["RequestRouter 统一鉴权", "按角色控制模块动作", "公开接口仅限排班查询和挂号创建", "写操作自动记录日志"], TEAL)
    add_card(s, 14.4, 2.05, 5.8, 4.7, "业务服务", ["Patient/Registration/Schedule", "Consultation/Examination", "Prescription/Inventory", "Billing/Statistics/Dashboard"], GREEN)
    add_card(s, 21.1, 2.05, 4.3, 4.7, "数据访问", ["QODBC/QMYSQL", "MySQL 表结构", "连接异常提示", "演示模式兜底"], ORANGE)

    # 10
    s = slides[9]
    add_title(s, "实施计划", "按软件工程过程推进，阶段产出清晰")
    periods = [
        ("需求与开题", "2025.12", "调研、模块划分、开题材料"),
        ("总体设计", "2026.01", "架构、数据库、原型"),
        ("核心开发", "2026.02-03", "患者、挂号、排班、接诊"),
        ("扩展开发", "2026.03-04", "处方、库存、收费、统计"),
        ("联调部署", "2026.04-05", "测试、麒麟/Windows 验证"),
        ("论文答辩", "2026.05-06", "论文、PPT、最终提交"),
    ]
    for i, (name, time, desc) in enumerate(periods):
        x = 1.0 + i * 4.15
        add_card(s, x, 2.25, 3.55, 4.2, name, [time, desc], [BLUE, TEAL, GREEN, ORANGE, RED, NAVY][i])

    # 11
    s = slides[10]
    add_title(s, "预期成果与特色", "答辩时重点突出“完整业务链”和“工程化实现”")
    add_card(s, 1.0, 2.15, 7.5, 4.75, "预期成果", ["可运行的 C/S 门诊管理系统", "MySQL 建库脚本和演示数据", "跨平台构建与运行脚本", "毕业论文、开题报告、答辩 PPT"], BLUE)
    add_card(s, 9.3, 2.15, 7.5, 4.75, "项目特色", ["患者入口与医院人员入口分离", "角色权限驱动菜单和接口", "排班号源影响预约挂号", "处方、药品、收费形成联动"], TEAL)
    add_card(s, 17.6, 2.15, 7.0, 4.75, "可展示亮点", ["候诊队列和叫号状态", "扫码入库和库存预警", "院长驾驶舱与费用统计", "操作日志和审计追踪"], GREEN)

    # 12
    s = slides[11]
    add_title(s, "风险分析与解决措施", "提前准备答辩老师可能关注的问题")
    add_card(s, 1.1, 2.1, 7.4, 4.85, "风险 1：数据库环境差异", ["措施：配置文件管理连接参数", "支持 QODBC/QMYSQL 两种方向", "保留演示模式验证界面和流程"], RED)
    add_card(s, 9.3, 2.1, 7.4, 4.85, "风险 2：业务状态不同步", ["措施：挂号成功扣减号源", "列表支持手动刷新和自动刷新", "核心状态由服务端统一维护"], ORANGE)
    add_card(s, 17.5, 2.1, 7.4, 4.85, "风险 3：模块多导致测试遗漏", ["措施：按角色设计测试用例", "从建档到收费做端到端验证", "操作日志辅助定位问题"], BLUE)

    # 13
    s = slides[12]
    thank = s.shapes.add_textbox(PCm(2.0), PCm(3.0), PCm(22.0), PCm(2.0))
    p = thank.text_frame.paragraphs[0]
    p.text = "请各位老师批评指正"
    p.font.name = "微软雅黑"
    p.font.size = PPt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    q = s.shapes.add_textbox(PCm(2.0), PCm(5.2), PCm(22.0), PCm(1.0))
    qp = q.text_frame.paragraphs[0]
    qp.text = "基于 Qt 和 MySQL 的医院门诊挂号与药品管理系统"
    qp.font.name = "微软雅黑"
    qp.font.size = PPt(17)
    qp.font.color.rgb = GRAY
    qp.alignment = PP_ALIGN.CENTER

    for i, slide in enumerate(slides, 1):
        add_footer(slide, i)

    prs.save(str(OUT_PPT))
    shutil.copy2(OUT_PPT, OUT_PPT_DOCS)


def build_notes():
    OUT_NOTES.write_text(
        """# 开题答辩讲稿提纲

1. 先说明背景：门诊业务高频、状态变化快，挂号、候诊、处方、药品、收费如果割裂，会带来重复录入和统计困难。
2. 再讲目标：本系统不是单一挂号页面，而是用 Qt/C++、TCP 服务端和 MySQL 做一个门诊基础业务闭环。
3. 功能结构重点讲三条链：患者挂号链、医生诊疗链、药品收费链；再补充权限、日志、统计这些管理能力。
4. 技术路线重点讲：客户端 ApiClient 发送 JSON Lines 请求，服务端 RequestRouter 鉴权并分发到业务 Service，最后由 DatabaseManager 访问 MySQL。
5. 数据库设计重点讲主外键关系：患者/医生/排班生成挂号，挂号关联病历、检查、处方，处方关联药品，账单和支付完成收费闭环。
6. 进度安排按需求、设计、核心开发、扩展开发、联调部署、论文答辩展开。
7. 预期成果强调可以运行、可演示、有数据库、有跨平台构建说明，并能在答辩现场展示主要流程。
""",
        encoding="utf-8",
    )


def main():
    draw_structure(OUT_STRUCTURE)
    draw_architecture(OUT_ARCH)
    draw_flow(OUT_FLOW)
    draw_database(OUT_DB)
    write_mermaid(OUT_MMD)
    build_report()
    build_ppt()
    build_notes()
    print(OUT_REPORT)
    print(OUT_PPT)
    print(OUT_STRUCTURE)
    print(OUT_NOTES)


if __name__ == "__main__":
    main()
